import tkinter as tk
from tkinter import ttk
from tkinter import filedialog
import os
from docx import Document
import PyPDF2
from pptx import Presentation
from openpyxl import load_workbook
import win32com.client as win32
import xlrd
import threading
import time


def search_keyword_in_ppt_file(file_path, keyword):
    presentation = Presentation(file_path)
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                if keyword in shape.text:
                    return True
    return False


def search_keyword_in_xls_file(file_path, keyword):
    # 打开Excel文件
    workbook = xlrd.open_workbook(file_path)

    # 获取所有工作表的名称
    sheet_names = workbook.sheet_names()

    for sheet_name in sheet_names:
        # 根据工作表名称选择工作表
        sheet = workbook.sheet_by_name(sheet_name)

        # 遍历单元格，检查关键词是否存在
        for row in range(sheet.nrows):
            for col in range(sheet.ncols):
                cell_value = sheet.cell_value(row, col)
                if keyword in str(cell_value):
                    # print(
                    #     f"关键词'{keyword}'存在于工作表'{sheet_name}'的单元格({row+1}, {col+1})")
                    # 如果你只想判断关键词是否存在，可直接返回True，避免继续遍历
                    return True

    # 如果关键词不存在于任何单元格
    # print(f"关键词'{keyword}'不存在于Excel文件中")
    return False


def search_keyword_in_xlsx_file(file_path, keyword):
    wb = load_workbook(filename=file_path)
    for sheet_name in wb.sheetnames:
        sheet = wb[sheet_name]
        for row in sheet.iter_rows():
            for cell in row:
                if cell.value and keyword in str(cell.value):
                    return True
    return False


def search_keyword_in_pdf_file(file_path, keyword):
    with open(file_path, 'rb') as f:
        pdf_reader = PyPDF2.PdfReader(f)
        for page in pdf_reader.pages:
            if keyword.lower() in page.extract_text().lower():
                return True
    return False


def search_keyword_in_doc_file(file_path, keyword):
    # 创建Word应用程序对象
    word_app = win32.Dispatch("Word.Application")
    # 打开Word文件
    doc = word_app.Documents.Open(file_path)
    # 将文档内容读取为纯文本
    doc_text = doc.Content.Text
    # 关闭Word文件
    doc.Close()
    # 退出Word应用程序
    word_app.Quit()

    # 检查关键词是否在文档内容中
    if keyword in doc_text:
        return True
    else:
        return False


def search_keyword_in_docx_file(file_path, keyword):
    doc = Document(file_path)
    for paragraph in doc.paragraphs:
        if keyword in paragraph.text:
            return True
    return False


def search_files(folder_path):
    file_extensions = ['.docx', '.doc', '.pdf', '.pptx', '.xlsx', '.xls']
    results = []
    for root, dirs, files in os.walk(folder_path):
        for file in files:
            filename, extension = os.path.splitext(file)
            if extension.lower() in file_extensions:
                file_path = os.path.join(root, file)
                results.append(file_path)
    return results


def on_select_folder():
    folder_path = filedialog.askdirectory()
    folder_label.config(text=folder_path)


def is_keyword_in_file(file_path, keyword):
    if file_path.endswith('.docx'):
        return search_keyword_in_docx_file(file_path, keyword)
    if file_path.endswith('.doc'):
        return search_keyword_in_doc_file(file_path, keyword)
    elif file_path.endswith('.pdf'):
        return search_keyword_in_pdf_file(file_path, keyword)
    elif file_path.endswith('.xls'):
        return search_keyword_in_xls_file(file_path, keyword)
    elif file_path.endswith('.xlsx'):
        return search_keyword_in_xlsx_file(file_path, keyword)
    elif file_path.endswith('.ppt'):
        return search_keyword_in_ppt_file(file_path, keyword)
    else:
        return False


def handle_note_text(text):
    note_text.config(text=text)
    # note_text.update()


output = ""
totalFilesCount = 0
handelFilesCount = 0
successFilesCount = 0
errorFilesCount = 0
errorTextFilePath = ''


def on_search():
    keyword = keyword_entry.get()
    folder = folder_label.cget("text")
    # result_text.config(text=f"关键字：{keyword}\n选择的文件夹：{folder}")
    global output
    global totalFilesCount
    global handelFilesCount
    global errorFilesCount
    global successFilesCount
    errorFilesCount = 0
    handelFilesCount = 0
    successFilesCount = 0
    output = ""

    # 调用函数搜索文件并打印路径
    files = search_files(folder)
    totalFilesCount = len(files)
    note_text.config(text=f"开始搜索：共{totalFilesCount}个文件")

    thread_loading_text = threading.Thread(target=handle_loading_text)
    thread_loading_text.start()

    for file in files:
        thread = threading.Thread(target=on_search_file, args=(file, keyword))
        thread.start()


def on_search_file(file, keyword):
    global output
    global handelFilesCount
    global errorFilesCount
    global successFilesCount
    try:

        if is_keyword_in_file(file, keyword):
            output += file + "\n"
            successFilesCount += 1
            print(file)
        handelFilesCount += 1
        # result_text.config(text=output)
        # 更新文本内容
        result_text.delete(1.0, tk.END)  # 清空现有内容
        result_text.insert(tk.END, output)

    except ZeroDivisionError as error:
        errorFilesCount += 1
        print(error)
        append_to_error_file(file)
        append_to_error_file(error)
        append_to_error_file("========================\n\n\n")


def handle_loading_text():
    while (totalFilesCount > handelFilesCount + errorFilesCount):
        time.sleep(0.1)
        if errorFilesCount > 0:
            note_text.config(
                text=f"正在搜索：共{totalFilesCount}个文件，已搜索{handelFilesCount}个文件,找到{successFilesCount}个文件,{errorFilesCount}个异常文件")
        else:
            note_text.config(
                text=f"正在搜索：共{totalFilesCount}个文件，已搜索{handelFilesCount}个文件,找到{successFilesCount}个文件")
    if errorFilesCount > 0:
        note_text.config(
            text=f"搜索结束：共搜索{totalFilesCount}个文件, 找到{successFilesCount}个文件,{errorFilesCount}个文件无法识别，请到error.txt查看")
    else:
        note_text.config(
            text=f"搜索结束：共搜索{totalFilesCount}个文件, 找到{successFilesCount}个文件")


def create_error_file(file_name):
    global errorTextFilePath
    errorTextFilePath = file_name
    # file_name = 'error.txt'
    try:
        # 检查文件是否存在，如果存在就删除
        if os.path.exists(file_name):
            os.remove(file_name)
            print(f"{file_name} 已存在，已删除")

        # 创建新文件
        with open(file_name, 'w') as file:
            file.write(time.localtime + "\n")
        print(f"{file_name} 创建成功")
    except Exception as e:
        print(f"发生错误: {e}")


def append_to_error_file(text):
    try:
        with open(errorTextFilePath, 'a') as file:  # 使用 'a' 模式打开文件，以追加模式写入文本
            file.write('无法识别文件： ' + text + '\n')  # 将文本写入文件末尾，并添加换行符
        print("文本成功添加到 error.txt 文件中")
    except Exception as e:
        print(f"发生错误: {e}")


def on_submit():
    # thread.start()
    # thread.join()
    print("开始搜索")
    keyword = keyword_entry.get()
    folder = folder_label.cget("text")
    create_error_file(os.path.join(folder, "error.txt"))
    if folder == "":
        handle_note_text("请选择文件夹！")
        return
    if keyword == "":
        handle_note_text("请输入关键字！")
        return
    note_text.config(text=f"搜索中。。。")
    # note_text.update()
    thread = threading.Thread(target=on_search)
    thread.start()


# 创建主窗口
window = tk.Tk()
window.title("搜索文件中的关键词")
window.geometry("600x520")  # 设置窗口大小
window['padx'] = 10
window['pady'] = 10

# 创建样式
style = ttk.Style()
style.configure("TLabel", background="white")
style.configure("TButton", padding=6, relief="flat")

# 创建关键字标签和输入框
keyword_label = ttk.Label(window, text="关键字：")
keyword_label.grid(row=0, column=0, padx=10, pady=10, sticky="e")
keyword_entry = ttk.Entry(window, width=30)
keyword_entry.grid(row=0, column=1, padx=10, pady=10, sticky="w")

# 创建选择文件夹按钮和标签
folder_button = ttk.Button(window, text="选择文件夹", command=on_select_folder)
folder_button.grid(row=1, column=0, padx=10, pady=10, sticky="e")
folder_label = ttk.Label(window, text="未选择文件夹")
folder_label.grid(row=1, column=1, padx=10, pady=10, sticky="w")

# 创建确定按钮
submit_button = ttk.Button(window, text="搜索", command=on_submit)
submit_button.grid(row=2, column=1, columnspan=2, padx=10, pady=10)

# loading
note_text = ttk.Label(window, text="")
note_text.grid(row=3, column=1, columnspan=2, padx=10, pady=10)

# 创建展示结果的 Text 组件和滚动条
result_text = tk.Text(window, wrap="word", height=20,
                      width=50, padx=10, pady=10)
result_text.grid(row=5, column=1)

scrollbar = tk.Scrollbar(window, command=result_text.yview)
scrollbar.grid(row=5, column=2, sticky="ns")
result_text.config(yscrollcommand=scrollbar.set)

# # 创建结果展示标签
# result_text = tk.Text(window, wrap="word", height=10, width=40)
# result_text.pack()

# # 创建滚动条
# scrollbar = tk.Scrollbar(root, command=result_text.yview)
# scrollbar.pack(side=tk.RIGHT, fill=tk.Y)
# result_text.config(yscrollcommand=scrollbar.set)
# result_text = ttk.Label(window, text="")
# result_text.grid(row=4, column=0, columnspan=2, padx=10, pady=10)
# result_text = tk.Text(window, width=40, height=10)
# result_text.grid(row=3, column=0, columnspan=2, padx=10, pady=10)
# result_text.configure(state="disabled") # 禁止编辑

# 运行主循环
window.mainloop()
